"""
PowerPoint document lifecycle agent.

Handles PPTX/PPT files through the full document lifecycle:
ingest → atomize → export, with support for slide masters,
speaker notes, images, charts, and PDF rendering via LibreOffice.

Requires ``python-pptx`` (pip install python-pptx).
"""

from __future__ import annotations

import logging
import uuid
from io import BytesIO
from typing import Any

from .base import (
    AgentCapability,
    AtomGroup,
    CanvasBundle,
    CanvasNode,
    ComplianceConstraints,
    DocumentAgent,
    ExportResult,
)

logger = logging.getLogger(__name__)

# ── Graceful import of python-pptx ────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    _HAS_PPTX = True
except ImportError:  # pragma: no cover
    _HAS_PPTX = False
    Presentation = None  # type: ignore[assignment,misc]

# ── Constants ──────────────────────────────────────────────────────────────
EMU_PER_INCH = 914400
EMU_PER_POINT = 914400 / 72  # 12700

# Placeholder type IDs for title detection (python-pptx placeholder_format.idx)
_TITLE_PLACEHOLDER_IDXS = {0, 1}  # 0 = title, 1 = center title

# Slide layout indices for a default presentation
_LAYOUT_TITLE_SLIDE = 0
_LAYOUT_TITLE_CONTENT = 1
_LAYOUT_BLANK = 5
_LAYOUT_TITLE_ONLY = 6


def _emu_to_points(emu: int) -> float:
    """Convert EMU (English Metric Units) to points."""
    return emu / EMU_PER_POINT


def _points_to_emu(points: float) -> int:
    """Convert points to EMU."""
    return int(points * EMU_PER_POINT)


def _is_title_placeholder(shape: Any) -> bool:
    """Return True if the shape is a title or center-title placeholder."""
    if not shape.has_text_frame:
        return False
    try:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
            pf_type = shape.placeholder_format.type
            if pf_type is not None:
                type_name = str(pf_type).upper()
                if "TITLE" in type_name or "CENTER_TITLE" in type_name:
                    return True
            # Fallback: idx 0 is conventionally the title
            idx = shape.placeholder_format.idx
            if idx == 0:
                return True
    except Exception:
        pass
    return False


def _is_body_placeholder(shape: Any) -> bool:
    """Return True if the shape is a body/content/subtitle placeholder (not title)."""
    if not shape.has_text_frame:
        return False
    try:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
            pf_type = shape.placeholder_format.type
            if pf_type is not None:
                type_name = str(pf_type).upper()
                # Exclude title placeholders
                if "TITLE" in type_name or "CENTER_TITLE" in type_name:
                    return False
                # Explicit body types
                if any(kw in type_name for kw in ("BODY", "SUBTITLE", "OBJECT")):
                    return True
            idx = shape.placeholder_format.idx
            # idx >= 1 (not title idx 0) is body/subtitle/content
            if idx >= 1:
                return True
    except Exception:
        pass
    return False


def _detect_aspect_ratio(slide_width: int, slide_height: int) -> str:
    """Determine format_type from slide dimensions in EMU."""
    if slide_height == 0:
        return "slide_16_9"
    ratio = slide_width / slide_height
    # 16:9 ≈ 1.778, 4:3 ≈ 1.333
    if abs(ratio - 16 / 9) < abs(ratio - 4 / 3):
        return "slide_16_9"
    return "slide_4_3"


def _make_id() -> str:
    """Generate a unique node id."""
    return str(uuid.uuid4())


class PptxAgent(DocumentAgent):
    """Document agent for Microsoft PowerPoint (.pptx / .ppt) files."""

    format_id = "pptx"
    display_name = "PowerPoint Document Agent"
    file_extensions = ["pptx", "ppt"]
    capabilities = {
        AgentCapability.READ_NATIVE,
        AgentCapability.WRITE_NATIVE,
        AgentCapability.SLIDE_MASTERS,
        AgentCapability.SPEAKER_NOTES,
        AgentCapability.IMAGES,
        AgentCapability.CHARTS,
        AgentCapability.PDF_RENDER,
    }

    def __init__(self) -> None:
        if not _HAS_PPTX:
            raise ImportError(
                "python-pptx is required for PptxAgent. "
                "Install it with: pip install python-pptx"
            )

    # ── INGEST ─────────────────────────────────────────────────────────

    async def ingest(self, file_bytes: bytes, filename: str) -> CanvasBundle:
        """Read a PPTX file and produce a CanvasBundle."""
        prs = Presentation(BytesIO(file_bytes))

        nodes: list[CanvasNode] = []
        image_count = 0
        chart_count = 0

        for slide_idx, slide in enumerate(prs.slides):
            # ── Page break between slides (not before the first) ──
            if slide_idx > 0:
                nodes.append(CanvasNode(
                    id=_make_id(),
                    type="page_break",
                    content={"slide_index": slide_idx},
                ))

            slide_notes = self._extract_notes(slide)
            slide_meta: dict[str, Any] = {"slide_index": slide_idx}
            if slide_notes:
                slide_meta["speaker_notes"] = slide_notes

            title_found = False

            for shape in slide.shapes:
                # ── Title placeholder ──
                if _is_title_placeholder(shape) and not title_found:
                    title_text = shape.text_frame.text.strip()
                    if title_text:
                        node_meta: dict[str, Any] = {"slide_index": slide_idx}
                        if slide_notes:
                            node_meta["speaker_notes"] = slide_notes
                        nodes.append(CanvasNode(
                            id=_make_id(),
                            type="heading",
                            content={"text": title_text, "level": 2},
                            provenance={
                                "source": "pptx_ingest",
                                "filename": filename,
                                "slide_index": slide_idx,
                                "element": "title",
                                **node_meta,
                            },
                        ))
                        title_found = True
                    continue

                # ── Body placeholder text ──
                if _is_body_placeholder(shape):
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if not para_text:
                            continue
                        nodes.append(CanvasNode(
                            id=_make_id(),
                            type="text_block",
                            content={"text": para_text},
                            provenance={
                                "source": "pptx_ingest",
                                "filename": filename,
                                "slide_index": slide_idx,
                                "element": "body",
                            },
                        ))
                    continue

                # ── Table ──
                if shape.has_table:
                    table = shape.table
                    rows: list[list[str]] = []
                    for row in table.rows:
                        rows.append([cell.text for cell in row.cells])
                    nodes.append(CanvasNode(
                        id=_make_id(),
                        type="table",
                        content={
                            "rows": rows,
                            "row_count": len(table.rows),
                            "col_count": len(table.columns),
                        },
                        provenance={
                            "source": "pptx_ingest",
                            "filename": filename,
                            "slide_index": slide_idx,
                            "element": "table",
                        },
                    ))
                    continue

                # ── Chart ──
                if shape.has_chart:
                    chart_count += 1
                    nodes.append(CanvasNode(
                        id=_make_id(),
                        type="text_block",
                        content={
                            "text": f"[Chart: {shape.chart.chart_title.text_frame.text if shape.chart.has_title else 'Untitled'}]",
                        },
                        provenance={
                            "source": "pptx_ingest",
                            "filename": filename,
                            "slide_index": slide_idx,
                            "element": "chart",
                        },
                    ))
                    continue

                # ── Image ──
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    img_name = ""
                    try:
                        img_name = shape.image.filename or ""
                    except Exception:
                        pass
                    nodes.append(CanvasNode(
                        id=_make_id(),
                        type="text_block",
                        content={"text": f"[Image: {img_name or 'embedded'}]"},
                        provenance={
                            "source": "pptx_ingest",
                            "filename": filename,
                            "slide_index": slide_idx,
                            "element": "image",
                            "image_filename": img_name,
                        },
                    ))
                    continue

                # ── Other text frame shapes (not placeholder) ──
                if shape.has_text_frame:
                    full_text = shape.text_frame.text.strip()
                    if full_text:
                        nodes.append(CanvasNode(
                            id=_make_id(),
                            type="text_block",
                            content={"text": full_text},
                            provenance={
                                "source": "pptx_ingest",
                                "filename": filename,
                                "slide_index": slide_idx,
                                "element": "text_box",
                            },
                        ))

        # ── Metadata from core properties ──
        meta: dict[str, Any] = {
            "filename": filename,
            "slide_count": len(prs.slides),
            "image_count": image_count,
            "chart_count": chart_count,
        }
        props = prs.core_properties
        if props.title:
            meta["title"] = props.title
        if props.author:
            meta["author"] = props.author
        if props.subject:
            meta["subject"] = props.subject
        if props.created:
            meta["created"] = props.created.isoformat()
        if props.modified:
            meta["modified"] = props.modified.isoformat()

        # ── Constraints from slide dimensions ──
        slide_width_pt = _emu_to_points(prs.slide_width)
        slide_height_pt = _emu_to_points(prs.slide_height)
        format_type = _detect_aspect_ratio(prs.slide_width, prs.slide_height)

        constraints = ComplianceConstraints(
            format_type=format_type,
            margins={
                "slide_width": slide_width_pt,
                "slide_height": slide_height_pt,
                "slide_width_emu": prs.slide_width,
                "slide_height_emu": prs.slide_height,
            },
        )

        return CanvasBundle(
            document_id=_make_id(),
            nodes=nodes,
            constraints=constraints,
            metadata=meta,
            source_format="pptx",
            source_agent=self.display_name,
        )

    # ── ATOMIZE ────────────────────────────────────────────────────────

    async def atomize(self, bundle: CanvasBundle) -> list[AtomGroup]:
        """Break a CanvasBundle into per-slide atom groups."""
        groups: list[AtomGroup] = []
        current_nodes: list[CanvasNode] = []
        current_heading: str | None = None
        slide_index = 0

        for node in bundle.nodes:
            if node.type == "page_break":
                # Flush the current group
                if current_nodes:
                    groups.append(AtomGroup(
                        nodes=list(current_nodes),
                        heading_text=current_heading,
                        suggested_category="slide",
                        suggested_tags=["presentation", f"slide_{slide_index + 1}"],
                        confidence=0.9,
                        source_slide=slide_index,
                    ))
                current_nodes = []
                current_heading = None
                slide_index += 1
                continue

            if node.type == "heading" and current_heading is None:
                current_heading = node.content.get("text", "")

            current_nodes.append(node)

        # Flush the last group
        if current_nodes:
            groups.append(AtomGroup(
                nodes=list(current_nodes),
                heading_text=current_heading,
                suggested_category="slide",
                suggested_tags=["presentation", f"slide_{slide_index + 1}"],
                confidence=0.9,
                source_slide=slide_index,
            ))

        return groups

    # ── EXPORT ─────────────────────────────────────────────────────────

    async def export(self, bundle: CanvasBundle) -> ExportResult:
        """Render a CanvasBundle to PPTX bytes."""
        prs = Presentation()

        # ── Apply slide dimensions from constraints ──
        margins = bundle.constraints.margins or {}
        if "slide_width_emu" in margins:
            prs.slide_width = int(margins["slide_width_emu"])
            prs.slide_height = int(margins["slide_height_emu"])
        elif bundle.constraints.format_type == "slide_4_3":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        else:
            # Default to 16:9
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        # ── Group nodes by page_break boundaries ──
        slide_groups = self._group_nodes_by_slide(bundle.nodes)

        warnings: list[str] = []
        slide_count = 0

        for group_nodes in slide_groups:
            slide_count += 1
            has_heading = any(n.type == "heading" for n in group_nodes)
            has_body = any(
                n.type in ("text_block", "bulleted_list", "numbered_list")
                for n in group_nodes
            )

            # ── Pick layout ──
            layout_idx = self._pick_layout(prs, has_heading, has_body)
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            # Track vertical position for manually-placed shapes
            body_top = Inches(1.8)
            body_left = Inches(0.5)
            body_width = prs.slide_width - Inches(1.0)
            current_top = body_top

            # Collect notes for this slide
            speaker_notes = ""

            for node in group_nodes:
                # ── Heading → title placeholder ──
                if node.type == "heading":
                    title_text = node.content.get("text", "")
                    self._set_title(slide, title_text, bundle.constraints)
                    # Grab speaker notes from node metadata/provenance
                    speaker_notes = node.provenance.get("speaker_notes", speaker_notes)
                    continue

                # ── Text block ──
                if node.type == "text_block":
                    text = node.content.get("text", "")
                    if not text:
                        continue
                    if has_heading and has_body and self._has_body_placeholder(slide):
                        self._add_to_body_placeholder(
                            slide, text, bundle.constraints
                        )
                    else:
                        current_top = self._add_textbox(
                            slide, text, body_left, current_top,
                            body_width, bundle.constraints,
                        )
                    continue

                # ── Bulleted list ──
                if node.type == "bulleted_list":
                    items = node.content.get("items", [])
                    if self._has_body_placeholder(slide):
                        for item in items:
                            self._add_to_body_placeholder(
                                slide, item, bundle.constraints, bulleted=True
                            )
                    else:
                        for item in items:
                            current_top = self._add_textbox(
                                slide, f"• {item}", body_left, current_top,
                                body_width, bundle.constraints,
                            )
                    continue

                # ── Numbered list ──
                if node.type == "numbered_list":
                    items = node.content.get("items", [])
                    if self._has_body_placeholder(slide):
                        for idx, item in enumerate(items, 1):
                            self._add_to_body_placeholder(
                                slide, f"{idx}. {item}", bundle.constraints
                            )
                    else:
                        for idx, item in enumerate(items, 1):
                            current_top = self._add_textbox(
                                slide, f"{idx}. {item}", body_left, current_top,
                                body_width, bundle.constraints,
                            )
                    continue

                # ── Table ──
                if node.type == "table":
                    rows_data = node.content.get("rows", [])
                    if rows_data:
                        try:
                            current_top = self._add_table(
                                slide, rows_data, body_left, current_top,
                                body_width, bundle.constraints,
                            )
                        except Exception as exc:
                            logger.error("PPTX_AGENT export table error: %s", exc)
                            warnings.append(f"Table render failed: {exc}")
                    continue

                # Check node provenance for speaker notes
                notes_text = node.provenance.get("speaker_notes", "")
                if notes_text and not speaker_notes:
                    speaker_notes = notes_text

            # ── Speaker notes ──
            if speaker_notes:
                self._set_speaker_notes(slide, speaker_notes)

        # ── Serialize ──
        buf = BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        out_filename = bundle.metadata.get("filename", "document.pptx")
        if not out_filename.lower().endswith((".pptx", ".ppt")):
            out_filename = out_filename.rsplit(".", 1)[0] + ".pptx"

        return ExportResult(
            file_bytes=pptx_bytes,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=out_filename,
            format="pptx",
            page_count=slide_count,
            warnings=warnings,
        )

    # ── Private helpers ────────────────────────────────────────────────

    @staticmethod
    def _extract_notes(slide: Any) -> str:
        """Safely extract speaker notes text from a slide."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    return notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        return ""

    @staticmethod
    def _group_nodes_by_slide(nodes: list[CanvasNode]) -> list[list[CanvasNode]]:
        """Split nodes into groups delimited by page_break nodes."""
        groups: list[list[CanvasNode]] = []
        current: list[CanvasNode] = []

        for node in nodes:
            if node.type == "page_break":
                if current:
                    groups.append(current)
                current = []
            else:
                current.append(node)

        if current:
            groups.append(current)

        return groups

    @staticmethod
    def _pick_layout(
        prs: Any, has_heading: bool, has_body: bool
    ) -> int:
        """Choose a slide layout index, falling back if index is out of range."""
        n_layouts = len(prs.slide_layouts)

        if has_heading and has_body:
            desired = _LAYOUT_TITLE_CONTENT  # 1
        elif has_heading:
            desired = _LAYOUT_TITLE_ONLY  # 6
        else:
            desired = _LAYOUT_BLANK  # 5

        if desired < n_layouts:
            return desired
        # Fallback: pick the last available layout
        return n_layouts - 1

    @staticmethod
    def _set_title(
        slide: Any, text: str, constraints: ComplianceConstraints
    ) -> None:
        """Set the title placeholder text on a slide."""
        if slide.shapes.title is not None:
            slide.shapes.title.text = text
            if constraints.font_family or constraints.font_size:
                for para in slide.shapes.title.text_frame.paragraphs:
                    for run in para.runs:
                        if constraints.font_family:
                            run.font.name = constraints.font_family
                        if constraints.font_size:
                            run.font.size = Pt(constraints.font_size)

    @staticmethod
    def _has_body_placeholder(slide: Any) -> bool:
        """Check if the slide has a body/content placeholder."""
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.idx >= 2:
                    return True
            except Exception:
                continue
        # Also check for idx 1 (subtitle/body in title+content layout)
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.idx == 1:
                    if shape.has_text_frame:
                        return True
            except Exception:
                continue
        return False

    @staticmethod
    def _add_to_body_placeholder(
        slide: Any,
        text: str,
        constraints: ComplianceConstraints,
        bulleted: bool = False,
    ) -> None:
        """Add text to the body/content placeholder on a slide."""
        body_shape = None
        # Prefer idx >= 2, then fall back to idx 1
        for shape in slide.placeholders:
            try:
                idx = shape.placeholder_format.idx
                if idx >= 2 and shape.has_text_frame:
                    body_shape = shape
                    break
            except Exception:
                continue

        if body_shape is None:
            for shape in slide.placeholders:
                try:
                    idx = shape.placeholder_format.idx
                    if idx == 1 and shape.has_text_frame:
                        body_shape = shape
                        break
                except Exception:
                    continue

        if body_shape is None:
            return

        tf = body_shape.text_frame

        # If placeholder already has non-empty text, add a new paragraph
        existing_text = tf.text.strip()
        if existing_text:
            para = tf.add_paragraph()
        else:
            # Use the first (empty) paragraph
            para = tf.paragraphs[0]

        para.text = text

        if bulleted:
            para.level = 0

        if constraints.font_family or constraints.font_size:
            for run in para.runs:
                if constraints.font_family:
                    run.font.name = constraints.font_family
                if constraints.font_size:
                    run.font.size = Pt(constraints.font_size)

    @staticmethod
    def _add_textbox(
        slide: Any,
        text: str,
        left: int,
        top: int,
        width: int,
        constraints: ComplianceConstraints,
    ) -> int:
        """Add a text box shape; returns the new top position after the box."""
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        para = tf.paragraphs[0]
        para.text = text

        if constraints.font_family or constraints.font_size:
            for run in para.runs:
                if constraints.font_family:
                    run.font.name = constraints.font_family
                if constraints.font_size:
                    run.font.size = Pt(constraints.font_size)

        return top + height + Inches(0.1)

    @staticmethod
    def _add_table(
        slide: Any,
        rows_data: list[list[str]],
        left: int,
        top: int,
        width: int,
        constraints: ComplianceConstraints,
    ) -> int:
        """Add a table shape; returns the new top position after the table."""
        n_rows = len(rows_data)
        n_cols = max(len(r) for r in rows_data) if rows_data else 1
        row_height = Inches(0.35)
        table_height = row_height * n_rows
        col_width = width // n_cols

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, left, top, width, table_height
        )
        table = table_shape.table

        for r_idx, row_data in enumerate(rows_data):
            for c_idx, cell_text in enumerate(row_data):
                if c_idx < n_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_text
                    if constraints.font_family or constraints.font_size:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if constraints.font_family:
                                    run.font.name = constraints.font_family
                                if constraints.font_size:
                                    run.font.size = Pt(constraints.font_size)

        return top + table_height + Inches(0.2)

    @staticmethod
    def _set_speaker_notes(slide: Any, notes_text: str) -> None:
        """Set speaker notes on a slide."""
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text
        except Exception as exc:
            logger.error("PPTX_AGENT failed to set speaker notes: %s", exc)
